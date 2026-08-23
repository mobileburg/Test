#!/usr/bin/env python3
"""Генерация двух версий презентации «Битриксуйте клиентские проекты»."""

from __future__ import annotations

import zipfile
from pathlib import Path

import qrcode
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "output"
LOGO = ASSETS / "bitrix24-logo-pack" / "bitrix24-logo" / "logo.png"
RUSSKIY_BIT_LOGO = ASSETS / "russkiy-bit-logo.png"
QR_PATH = ASSETS / "vk-anchebykin-qr.png"
VK_URL = "https://vk.ru/anchebykin"

CYAN = "2FC7F7"
BLUE = "409EEF"
NAVY = "0066A1"
LIME = "BDF300"
INK = "10293E"
MUTED = "5A7082"
PALE = "F2F8FC"
LINE = "D8E9F3"
WHITE = "FFFFFF"
FONT_HEAD = "Montserrat"
FONT_BODY = "Open Sans"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_shape_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_shape_line(shape, color: str, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = INK,
    font: str = FONT_BODY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
    line_spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = LINE,
    radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(
        radius_shape, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_shape_fill(shape, fill)
    set_shape_line(shape, line, 1.1)
    return shape


def add_label(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    color: str = INK,
    size: float = 15,
    line: str = LINE,
):
    shape = add_round_rect(slide, x, y, w, h, fill, line)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.12)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_HEAD
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    return shape


def add_double_arrow(slide, x1: float, y1: float, x2: float, y2: float, width=1.6):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_shape_line(line, BLUE, width)
    line_el = line._element.spPr.ln
    head = OxmlElement("a:headEnd")
    head.set("type", "triangle")
    head.set("w", "sm")
    head.set("len", "sm")
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line_el.append(head)
    line_el.append(tail)
    return line


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, width=2.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_shape_line(line, NAVY, width)
    line_el = line._element.spPr.ln
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    line_el.append(tail)
    return line


def add_base(slide, width: float, height: float, _number: int) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)

    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(width), Inches(0.08)
    )
    set_shape_fill(top, CYAN)
    top.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(width - (1.45 if width < 7 else 1.8)),
        Inches(height - (1.25 if width < 7 else 1.45)),
        Inches(2.1 if width < 7 else 2.5),
        Inches(2.1 if width < 7 else 2.5),
    )
    set_shape_fill(circle, CYAN, 88)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.8),
        Inches(height - 1.15),
        Inches(1.65),
        Inches(1.65),
    )
    set_shape_fill(circle2, LIME, 82)
    circle2.line.fill.background()

    logo_width = 1.35 if width < 7 else 1.75
    slide.shapes.add_picture(
        str(RUSSKIY_BIT_LOGO),
        Inches(width - logo_width - (0.4 if width < 7 else 0.72)),
        Inches(0.28 if width < 7 else 0.22),
        width=Inches(logo_width),
    )


def add_heading(
    slide,
    title: str,
    width: float,
    *,
    mobile: bool,
    subtitle: str | None = None,
) -> float:
    x = 0.52 if mobile else 0.72
    y = 0.55 if mobile else 0.48
    title_size = 21.5 if mobile else 27
    title_h = 1.15 if mobile else 0.65
    title_w = width - x - (1.95 if mobile else 2.65)
    add_text(
        slide,
        title,
        x,
        y,
        title_w,
        title_h,
        size=title_size,
        color=INK,
        font=FONT_HEAD,
        bold=True,
        line_spacing=0.9,
    )
    if subtitle:
        sub_y = y + (1.0 if mobile else 0.64)
        add_text(
            slide,
            subtitle,
            x,
            sub_y,
            width - 2 * x,
            0.72 if mobile else 0.45,
            size=11.5 if mobile else 12.5,
            color=MUTED,
            font=FONT_BODY,
            line_spacing=0.95,
        )
        return sub_y + (0.82 if mobile else 0.52)
    return y + (1.18 if mobile else 0.78)


def add_logo_block(slide, x: float, y: float, w: float, h: float):
    halo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - 0.22),
        Inches(y - 0.22),
        Inches(w + 0.44),
        Inches(h + 0.44),
    )
    set_shape_fill(halo, CYAN, 84)
    halo.line.fill.background()
    card = add_round_rect(slide, x, y, w, h, WHITE, CYAN)
    set_shape_line(card, CYAN, 1.8)
    logo_w = w * 0.77
    logo_h = logo_w * 153 / 760
    slide.shapes.add_picture(
        str(LOGO),
        Inches(x + (w - logo_w) / 2),
        Inches(y + (h - logo_h) / 2),
        width=Inches(logo_w),
        height=Inches(logo_h),
    )


def slide_one(prs: Presentation, width: float, height: float, mobile: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, width, height, 1)
    content_top = add_heading(
        slide, "Битриксуйте клиентские проекты", width, mobile=mobile
    )

    if mobile:
        cx, cy, cw, ch = 1.17, 4.18, 3.28, 1.22
        nodes = {
            "HR": (1.84, content_top + 0.18, 1.94, 0.66),
            "Маркетинг": (0.48, 2.84, 2.15, 0.76),
            "Продажи": (3.01, 2.84, 2.15, 0.76),
            "Финансы": (0.48, 6.08, 2.15, 0.76),
            "Digital": (3.01, 6.08, 2.15, 0.76),
        }
        arrow_points = {
            "HR": (2.81, content_top + 0.84, 2.81, cy),
            "Маркетинг": (2.63, 3.22, cx, cy + 0.35),
            "Продажи": (3.01, 3.22, cx + cw, cy + 0.35),
            "Финансы": (2.63, 6.08, cx, cy + ch - 0.25),
            "Digital": (3.01, 6.08, cx + cw, cy + ch - 0.25),
        }
        node_size = 13.5
    else:
        cx, cy, cw, ch = 5.0, 3.1, 3.33, 1.24
        nodes = {
            "HR": (5.65, content_top + 0.08, 2.03, 0.66),
            "Маркетинг": (1.0, 2.6, 2.38, 0.78),
            "Продажи": (9.95, 2.6, 2.38, 0.78),
            "Финансы": (2.25, 5.0, 2.38, 0.78),
            "Digital": (8.7, 5.0, 2.38, 0.78),
        }
        arrow_points = {
            "HR": (6.67, content_top + 0.74, 6.67, cy),
            "Маркетинг": (3.38, 2.99, cx, cy + 0.32),
            "Продажи": (9.95, 2.99, cx + cw, cy + 0.32),
            "Финансы": (4.63, 5.0, cx + 0.46, cy + ch),
            "Digital": (8.7, 5.0, cx + cw - 0.46, cy + ch),
        }
        node_size = 15

    for points in arrow_points.values():
        add_double_arrow(slide, *points)
    add_logo_block(slide, cx, cy, cw, ch)
    for text, (x, y, w, h) in nodes.items():
        add_label(
            slide,
            text,
            x,
            y,
            w,
            h,
            fill=PALE,
            color=NAVY,
            size=node_size,
            line=LINE,
        )


def slide_two(prs: Presentation, width: float, height: float, mobile: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, width, height, 2)
    top = add_heading(
        slide,
        "Об этом не говорят в рекламе",
        width,
        mobile=mobile,
        subtitle="С этим столкнётся ваш клиент — или почему проваливаются внедрения",
    )
    items = [
        "Автоматизация хаоса",
        "Лоскутная автоматизация",
        "Техничка вместо сценариев",
        "Трудности перевода",
        "Саботаж персонала",
        "Отсутствие драйвера перемен",
    ]

    if mobile:
        x_positions = [0.46, 2.96]
        y0, card_w, card_h, gap_y = top + 0.05, 2.28, 0.86, 0.25
        centers = []
        cards = []
        for i, item in enumerate(items):
            row, col = divmod(i, 2)
            x = x_positions[col]
            y = y0 + row * (card_h + gap_y)
            cards.append(
                (
                    item,
                    x,
                    y,
                    card_w,
                    card_h,
                    PALE if i % 2 == 0 else WHITE,
                    11.2,
                    CYAN if i % 2 == 0 else LINE,
                )
            )
            centers.append((x + card_w / 2, y + card_h))
        rule_y = y0 + 3 * (card_h + gap_y) + 0.15
        result_y = rule_y + 0.75
        result_x, result_w, result_h = 0.46, 4.78, 1.15
    else:
        x_positions = [0.72, 4.67, 8.62]
        y0, card_w, card_h, gap_y = top + 0.13, 3.45, 0.77, 0.28
        centers = []
        cards = []
        for i, item in enumerate(items):
            row, col = divmod(i, 3)
            x = x_positions[col]
            y = y0 + row * (card_h + gap_y)
            cards.append(
                (
                    item,
                    x,
                    y,
                    card_w,
                    card_h,
                    PALE if i % 2 == 0 else WHITE,
                    12.4,
                    CYAN if i % 2 == 0 else LINE,
                )
            )
            centers.append((x + card_w / 2, y + card_h))
        rule_y = y0 + 2 * (card_h + gap_y) + 0.14
        result_y = rule_y + 0.65
        result_x, result_w, result_h = 3.17, 7.0, 0.88

    rule = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(min(x for x, _ in centers)),
        Inches(rule_y),
        Inches(max(x for x, _ in centers)),
        Inches(rule_y),
    )
    set_shape_line(rule, BLUE, 1.3, 22)
    for cx, cy in centers:
        stem = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(cx),
            Inches(cy + 0.07),
            Inches(cx),
            Inches(rule_y),
        )
        set_shape_line(stem, BLUE, 1.15, 25)

    mid_x = width / 2
    add_arrow(slide, mid_x, rule_y, mid_x, result_y - 0.08, 2.4)
    for item, x, y, card_w, card_h, fill, size, line in cards:
        add_label(
            slide,
            item,
            x,
            y,
            card_w,
            card_h,
            fill=fill,
            color=INK,
            size=size,
            line=line,
        )
    result = add_round_rect(
        slide, result_x, result_y, result_w, result_h, NAVY, NAVY
    )
    frame = result.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.2)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Непонимание, страх неудачи, самосаботаж"
    run.font.name = FONT_HEAD
    run.font.size = Pt(14 if mobile else 16)
    run.font.bold = True
    run.font.color.rgb = rgb(WHITE)


def add_point(slide, number: str, text: str, x: float, y: float, w: float, mobile: bool):
    card = add_round_rect(slide, x, y, w, 0.82 if mobile else 0.72, WHITE, LINE)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.15),
        Inches(y + (0.13 if mobile else 0.1)),
        Inches(0.56 if mobile else 0.5),
        Inches(0.56 if mobile else 0.5),
    )
    set_shape_fill(dot, CYAN)
    dot.line.fill.background()
    add_text(
        slide,
        number,
        x + 0.15,
        y + (0.13 if mobile else 0.1),
        0.56 if mobile else 0.5,
        0.56 if mobile else 0.5,
        size=10,
        color=NAVY,
        font=FONT_HEAD,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        text,
        x + (0.86 if mobile else 0.78),
        y + 0.1,
        w - (1.02 if mobile else 0.94),
        0.62 if mobile else 0.52,
        size=11.5 if mobile else 12.4,
        color=INK,
        font=FONT_HEAD,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return card


def slide_three(prs: Presentation, width: float, height: float, mobile: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, width, height, 3)
    top = add_heading(
        slide,
        "Вывод: нужен честный грамотный партнёр",
        width,
        mobile=mobile,
    )
    points = [
        "Говорит на одном языке с вами и клиентом",
        "Создаст систему на основе понимания бизнеса",
        "Обучит и поддержит",
    ]
    if mobile:
        x, card_w = 0.46, 4.78
        y = top + 0.08
        for i, text in enumerate(points, 1):
            add_point(slide, f"0{i}", text, x, y, card_w, True)
            y += 1.0
        cta_y, cta_h = y + 0.15, 3.75
        cta = add_round_rect(slide, x, cta_y, card_w, cta_h, NAVY, NAVY)
        add_text(
            slide,
            "Пишите мне в ВК",
            x + 0.32,
            cta_y + 0.32,
            card_w - 0.64,
            0.42,
            size=18,
            color=WHITE,
            font=FONT_HEAD,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "и получите чек-лист\n«10 ошибок внедрения Битрикс24.\nКак их избежать»",
            x + 0.32,
            cta_y + 0.82,
            card_w - 0.64,
            1.12,
            size=11.5,
            color=WHITE,
            font=FONT_BODY,
            align=PP_ALIGN.CENTER,
            line_spacing=0.95,
        )
        qr_size = 1.46
        qr = slide.shapes.add_picture(
            str(QR_PATH),
            Inches(x + (card_w - qr_size) / 2),
            Inches(cta_y + 1.98),
            width=Inches(qr_size),
            height=Inches(qr_size),
        )
        qr.click_action.hyperlink.address = VK_URL
        add_text(
            slide,
            "vk.ru/anchebykin",
            x + 0.4,
            cta_y + 3.48,
            card_w - 0.8,
            0.18,
            size=8.5,
            color=CYAN,
            font=FONT_BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    else:
        left_x, left_w = 0.72, 6.15
        y = top + 0.2
        for i, text in enumerate(points, 1):
            add_point(slide, f"0{i}", text, left_x, y, left_w, False)
            y += 0.92
        cta_x, cta_y, cta_w, cta_h = 7.34, top + 0.2, 5.28, 3.43
        add_round_rect(slide, cta_x, cta_y, cta_w, cta_h, NAVY, NAVY)
        add_text(
            slide,
            "Пишите мне в ВК",
            cta_x + 0.38,
            cta_y + 0.34,
            3.0,
            0.42,
            size=19,
            color=WHITE,
            font=FONT_HEAD,
            bold=True,
        )
        add_text(
            slide,
            "и получите чек-лист\n«10 ошибок внедрения Битрикс24.\nКак их избежать»",
            cta_x + 0.38,
            cta_y + 0.92,
            3.1,
            1.2,
            size=11.8,
            color=WHITE,
            font=FONT_BODY,
            line_spacing=0.95,
        )
        qr_size = 1.62
        qr = slide.shapes.add_picture(
            str(QR_PATH),
            Inches(cta_x + cta_w - qr_size - 0.38),
            Inches(cta_y + 0.44),
            width=Inches(qr_size),
            height=Inches(qr_size),
        )
        qr.click_action.hyperlink.address = VK_URL
        link_shape = add_label(
            slide,
            "vk.ru/anchebykin  →",
            cta_x + 0.38,
            cta_y + 2.56,
            cta_w - 0.76,
            0.52,
            fill=LIME,
            color=INK,
            size=10.5,
            line=LIME,
        )
        link_shape.click_action.hyperlink.address = VK_URL


def make_qr() -> None:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=14,
        border=4,
    )
    qr.add_data(VK_URL)
    qr.make(fit=True)
    image = qr.make_image(fill_color=f"#{INK}", back_color=f"#{WHITE}").convert("RGB")
    draw = ImageDraw.Draw(image)
    side = image.width
    badge = int(side * 0.19)
    left = (side - badge) // 2
    top = (side - badge) // 2
    pad = max(7, int(side * 0.013))
    draw.rounded_rectangle(
        (left - pad, top - pad, left + badge + pad, top + badge + pad),
        radius=int(badge * 0.22),
        fill=f"#{WHITE}",
    )
    draw.rounded_rectangle(
        (left, top, left + badge, top + badge),
        radius=int(badge * 0.18),
        fill="#0077FF",
    )
    font_candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf",
    ]
    font_path = next((path for path in font_candidates if Path(path).exists()), None)
    font = ImageFont.truetype(font_path, int(badge * 0.39)) if font_path else None
    text = "VK"
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(
        ((side - tw) / 2, (side - th) / 2 - bbox[1]),
        text,
        fill=f"#{WHITE}",
        font=font,
    )
    image.save(QR_PATH, quality=100)


def build(filename: str, width: float, height: float, mobile: bool) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    prs.core_properties.title = "Битриксуйте клиентские проекты"
    prs.core_properties.subject = "Битрикс24: клиентские проекты и внедрение"
    prs.core_properties.author = "Презентация для Битрикс24"
    prs.core_properties.language = "ru-RU"
    slide_one(prs, width, height, mobile)
    slide_two(prs, width, height, mobile)
    slide_three(prs, width, height, mobile)
    path = OUTPUT / filename
    prs.save(path)
    return path


def main() -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    QR_PATH.parent.mkdir(parents=True, exist_ok=True)
    make_qr()
    desktop = build("Bitrix24_klientskie_proekty_PC_16x9.pptx", 13.333, 7.5, False)
    mobile = build("Bitrix24_klientskie_proekty_Mobile_9x16.pptx", 5.625, 10.0, True)
    for path in (desktop, mobile):
        print(f"Создан: {path.relative_to(ROOT)}")

    archive = OUTPUT / "Bitrix24_klientskie_proekty_materialy.zip"
    if archive.exists():
        archive.unlink()
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as bundle:
        for path in sorted(OUTPUT.iterdir()):
            if path.is_file() and path != archive:
                bundle.write(path, arcname=path.name)
    print(f"Создан: {archive.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
